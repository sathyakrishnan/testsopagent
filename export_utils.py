"""
Kevin AI - Export Utilities
Generate PDF and PPTX reports from results

Version: 2.0
Date: January 6, 2026
"""

from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_LEFT
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
import json
import base64
import subprocess
import tempfile
import os


def mermaid_to_image(mermaid_code, output_path, format='png'):
    """
    Convert Mermaid diagram to PNG or JPEG using mermaid.ink API
    
    Args:
        mermaid_code: Mermaid diagram code
        output_path: Output file path
        format: 'png' or 'jpg'
    """
    try:
        # Use mermaid.ink free API service
        import requests
        
        # Encode mermaid code to base64
        mermaid_bytes = mermaid_code.encode('utf-8')
        base64_bytes = base64.b64encode(mermaid_bytes)
        base64_string = base64_bytes.decode('utf-8')
        
        # Use mermaid.ink API
        if format.lower() == 'jpg' or format.lower() == 'jpeg':
            url = f"https://mermaid.ink/img/{base64_string}?type=jpeg"
        else:
            url = f"https://mermaid.ink/img/{base64_string}?type=png"
        
        # Download image
        response = requests.get(url, timeout=30)
        response.raise_for_status()
        
        # Save to file
        with open(output_path, 'wb') as f:
            f.write(response.content)
        
        return output_path
        
    except Exception as e:
        print(f"Error converting Mermaid to image: {e}")
        # Fallback: Create a simple text-based diagram placeholder
        from PIL import Image, ImageDraw, ImageFont
        
        img = Image.new('RGB', (800, 600), color='white')
        draw = ImageDraw.Draw(img)
        
        # Draw text
        text = "Mermaid Diagram\n(View full diagram in outputs)"
        draw.text((50, 50), text, fill='black')
        
        img.save(output_path)
        return output_path


def create_pdf_report(result, output_path):
    """
    Create comprehensive PDF report with all artifacts
    """
    doc = SimpleDocTemplate(output_path, pagesize=letter)
    story = []
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1f77b4'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#2c3e50'),
        spaceAfter=12,
        spaceBefore=12
    )
    
    # Title Page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("Kevin AI", title_style))
    story.append(Paragraph("SOP Automation Analysis Report", styles['Heading2']))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph(f"Session ID: {result['session_id']}", styles['Normal']))
    story.append(Paragraph(f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", styles['Normal']))
    story.append(Paragraph(f"Domain: {result.get('domain', 'N/A').capitalize()}", styles['Normal']))
    story.append(PageBreak())
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", heading_style))
    
    automation_opps = result.get('automation_opportunities', [])
    test_cases = result.get('test_cases', [])
    total_savings = sum(opp.get('estimated_savings_annual', 0) for opp in automation_opps)
    
    summary_data = [
        ['Metric', 'Value'],
        ['Automation Opportunities', str(len(automation_opps))],
        ['Test Cases Generated', str(len(test_cases))],
        ['Estimated Annual Savings', f'${total_savings:,.0f}'],
        ['Processing Time', '~14 minutes'],
    ]
    
    summary_table = Table(summary_data, colWidths=[3*inch, 3*inch])
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    story.append(summary_table)
    story.append(Spacer(1, 0.3*inch))
    story.append(PageBreak())
    
    # Current State Process Map
    story.append(Paragraph("Current State Process Map", heading_style))
    if result.get('current_state_map'):
        story.append(Paragraph("Mermaid Diagram:", styles['Normal']))
        story.append(Spacer(1, 0.1*inch))
        # Truncate if too long
        map_text = result['current_state_map'][:2000] + "..." if len(result['current_state_map']) > 2000 else result['current_state_map']
        story.append(Paragraph(f"<pre>{map_text}</pre>", styles['Code']))
    else:
        story.append(Paragraph("No current state map available.", styles['Normal']))
    story.append(PageBreak())
    
    # Future State Process Map
    story.append(Paragraph("Future State Process Map", heading_style))
    if result.get('future_state_map'):
        story.append(Paragraph("Optimized Mermaid Diagram:", styles['Normal']))
        story.append(Spacer(1, 0.1*inch))
        map_text = result['future_state_map'][:2000] + "..." if len(result['future_state_map']) > 2000 else result['future_state_map']
        story.append(Paragraph(f"<pre>{map_text}</pre>", styles['Code']))
    else:
        story.append(Paragraph("No future state map available.", styles['Normal']))
    story.append(PageBreak())
    
    # Automation Opportunities
    story.append(Paragraph("Automation Opportunities Matrix", heading_style))
    if automation_opps:
        # Create table data
        table_data = [
            ['ID', 'Stage', 'Activity', 'Type', 'Impact', 'Savings/Year', 'Priority']
        ]
        
        for i, opp in enumerate(automation_opps, 1):
            table_data.append([
                f"AO-{i:02d}",
                str(opp.get('step_id', 'N/A'))[:15],
                str(opp.get('description', 'N/A'))[:40] + "...",
                str(opp.get('automation_type', 'N/A'))[:15],
                f"{opp.get('time_savings_hours_per_week', 0) * 2}%",
                f"${opp.get('estimated_savings_annual', 0):,.0f}",
                str(opp.get('priority', 'P2'))
            ])
        
        opp_table = Table(table_data, colWidths=[0.5*inch, 0.8*inch, 2*inch, 1*inch, 0.6*inch, 1*inch, 0.6*inch])
        opp_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))
        
        story.append(opp_table)
        story.append(Spacer(1, 0.3*inch))
        
        # Summary by type
        story.append(Paragraph("Automation by Type:", styles['Normal']))
        automation_types = {}
        for opp in automation_opps:
            atype = opp.get('automation_type', 'Unknown')
            automation_types[atype] = automation_types.get(atype, 0) + 1
        
        for atype, count in automation_types.items():
            story.append(Paragraph(f"• {atype}: {count} opportunities", styles['Normal']))
        
    else:
        story.append(Paragraph("No automation opportunities identified.", styles['Normal']))
    story.append(PageBreak())
    
    # Test Cases Summary
    story.append(Paragraph("Test Cases", heading_style))
    if test_cases:
        # Test cases table
        table_data = [
            ['ID', 'Test Case', 'Type', 'Priority', 'Status']
        ]
        
        for i, test in enumerate(test_cases[:20], 1):  # First 20
            table_data.append([
                f"TC-{i:02d}",
                str(test.get('test_name', 'N/A'))[:50],
                str(test.get('test_type', 'Functional')),
                str(test.get('priority', 'Medium')),
                'Auto-Ready' if test.get('automated', True) else 'Manual'
            ])
        
        test_table = Table(table_data, colWidths=[0.6*inch, 3*inch, 1*inch, 0.8*inch, 1*inch])
        test_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ]))
        
        story.append(test_table)
        story.append(Spacer(1, 0.3*inch))
        
        # Detailed test cases (first 3)
        story.append(PageBreak())
        story.append(Paragraph("Detailed Test Case Examples", heading_style))
        
        for i, test in enumerate(test_cases[:3], 1):
            story.append(Paragraph(f"<b>TC-{i:02d}: {test.get('test_name', 'N/A')}</b>", styles['Heading3']))
            story.append(Paragraph(f"<b>Type:</b> {test.get('test_type', 'Functional')}", styles['Normal']))
            story.append(Paragraph(f"<b>Priority:</b> {test.get('priority', 'Medium')}", styles['Normal']))
            story.append(Paragraph(f"<b>Process Step:</b> {test.get('process_step', 'N/A')}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
            
            story.append(Paragraph("<b>Description:</b>", styles['Normal']))
            story.append(Paragraph(test.get('description', 'N/A'), styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
            
            story.append(Paragraph("<b>Pre-conditions:</b>", styles['Normal']))
            preconditions = test.get('preconditions', ['System is available'])
            for pc in preconditions:
                story.append(Paragraph(f"• {pc}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
            
            story.append(Paragraph("<b>Test Steps:</b>", styles['Normal']))
            test_steps = test.get('test_steps', ['Execute test'])
            for j, step in enumerate(test_steps, 1):
                story.append(Paragraph(f"{j}. {step}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
            
            story.append(Paragraph("<b>Expected Result:</b>", styles['Normal']))
            story.append(Paragraph(test.get('expected_result', 'Test passes'), styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
        
        # Summary by type
        story.append(Spacer(1, 0.2*inch))
        story.append(Paragraph("Test Coverage Summary:", styles['Normal']))
        test_types = {}
        for test in test_cases:
            ttype = test.get('test_type', 'Unknown')
            test_types[ttype] = test_types.get(ttype, 0) + 1
        
        for ttype, count in test_types.items():
            story.append(Paragraph(f"• {ttype}: {count} tests", styles['Normal']))
        
    else:
        story.append(Paragraph("No test cases generated.", styles['Normal']))
    
    # Build PDF
    doc.build(story)
    return output_path


def create_pptx_report(result, output_path):
    """
    Create PowerPoint presentation with key findings
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Kevin AI"
    subtitle.text = f"SOP Automation Analysis Report\nSession: {result['session_id']}\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Executive Summary Slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    
    tf = body_shape.text_frame
    automation_opps = result.get('automation_opportunities', [])
    test_cases = result.get('test_cases', [])
    total_savings = sum(opp.get('estimated_savings_annual', 0) for opp in automation_opps)
    
    tf.text = f"Automation Opportunities: {len(automation_opps)}"
    
    p = tf.add_paragraph()
    p.text = f"Test Cases Generated: {len(test_cases)}"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"Estimated Annual Savings: ${total_savings:,.0f}"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"Processing Time: ~14 minutes"
    p.level = 0
    
    # Automation Opportunities Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "Top Automation Opportunities"
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, opp in enumerate(automation_opps[:5], 1):
        p = tf.add_paragraph()
        p.text = f"{opp.get('description', 'N/A')} - ${opp.get('estimated_savings_annual', 0):,.0f}/year"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"Type: {opp.get('automation_type', 'N/A')} | Complexity: {opp.get('complexity', 'N/A')}"
        p.level = 1
    
    # Test Cases Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "Test Coverage"
    
    tf = body_shape.text_frame
    tf.text = f"Total Test Cases: {len(test_cases)}"
    
    # Count test types
    test_types = {}
    for test in test_cases:
        test_type = test.get('test_type', 'Unknown')
        test_types[test_type] = test_types.get(test_type, 0) + 1
    
    for test_type, count in test_types.items():
        p = tf.add_paragraph()
        p.text = f"{test_type}: {count} tests"
        p.level = 1
    
    # Next Steps Slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = "Recommended Next Steps"
    
    tf = body_shape.text_frame
    tf.text = "Review automation opportunities with stakeholders"
    
    p = tf.add_paragraph()
    p.text = "Prioritize implementations based on ROI"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Execute test cases for validation"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Deploy generated code to production"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Monitor KPIs and iterate"
    p.level = 0
    
    # Save presentation
    prs.save(output_path)
    return output_path
